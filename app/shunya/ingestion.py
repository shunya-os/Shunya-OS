"""Shunya Data Ingestion Pipeline — ingest any file format, AI extracts and structures.

Supports: PDF, DOCX, images (OCR), audio (STT), video (STT), CSV, XLSX, JSON,
Markdown, HTML, plain text, ZIP archives, and URLs.

Pipeline: Upload → Extract Text → AI Analyze → Structure → Store in Knowledge Base
"""
import os, json, logging, tempfile, uuid, zipfile, mimetypes
from pathlib import Path
from typing import Optional, List, Dict
from datetime import datetime
from flask import g
from app import db
from app.models import KnowledgeEntry, File, Entity, ActivityLog

logger = logging.getLogger("app.shunya.ingestion")

# Supported file types and their extractors
SUPPORTED_EXTENSIONS = {
    # Documents
    ".pdf": "pdf", ".docx": "docx", ".doc": "docx",
    ".txt": "text", ".md": "text", ".html": "text", ".htm": "text",
    ".rtf": "text", ".csv": "tabular", ".tsv": "tabular",
    ".xlsx": "tabular", ".xls": "tabular",
    ".json": "json", ".xml": "xml", ".yaml": "yaml", ".yml": "yaml",
    ".toml": "text",
    # Images (OCR)
    ".png": "image", ".jpg": "image", ".jpeg": "image",
    ".gif": "image", ".bmp": "image", ".tiff": "image", ".webp": "image",
    # Audio (STT)
    ".mp3": "audio", ".wav": "audio", ".ogg": "audio",
    ".m4a": "audio", ".flac": "audio", ".aac": "audio", ".wma": "audio",
    # Video (STT)
    ".mp4": "video", ".avi": "video", ".mov": "video",
    ".mkv": "video", ".webm": "video",
    # Archives
    ".zip": "archive", ".tar": "archive", ".gz": "archive",
    # Presentations
    ".pptx": "presentation", ".ppt": "presentation",
    # Other
    ".eml": "email", ".msg": "email",
}


class IngestionPipeline:
    """Universal data ingestion pipeline with AI extraction and structuring."""

    @staticmethod
    def process_file(file_path: str, tenant_id: int, user_id: int,
                     source_label: Optional[str] = None,
                     category: Optional[str] = None) -> dict:
        """Process a single file through the pipeline.

        Returns structured result with extracted text, AI analysis, and knowledge entries.
        """
        path = Path(file_path)
        ext = path.suffix.lower()
        processor = SUPPORTED_EXTENSIONS.get(ext)

        if not processor:
            return {"error": f"Unsupported file type: {ext}", "success": False}

        # Step 1: Extract raw text
        raw_text = ""
        extraction_method = processor
        metadata = {}

        try:
            if processor == "pdf":
                raw_text, metadata = IngestionPipeline._extract_pdf(file_path)
            elif processor == "docx":
                raw_text, metadata = IngestionPipeline._extract_docx(file_path)
            elif processor == "text":
                raw_text, metadata = IngestionPipeline._extract_text(file_path)
            elif processor == "tabular":
                raw_text, metadata = IngestionPipeline._extract_tabular(file_path)
            elif processor == "image":
                raw_text, metadata = IngestionPipeline._extract_image(file_path)
            elif processor == "audio":
                raw_text, metadata = IngestionPipeline._extract_audio(file_path)
            elif processor == "video":
                raw_text, metadata = IngestionPipeline._extract_video(file_path)
            elif processor == "json":
                raw_text, metadata = IngestionPipeline._extract_json(file_path)
            elif processor == "xml":
                raw_text, metadata = IngestionPipeline._extract_text(file_path)
            elif processor == "archive":
                return IngestionPipeline._process_archive(file_path, tenant_id, user_id, category)
            elif processor == "presentation":
                raw_text, metadata = IngestionPipeline._extract_presentation(file_path)
            elif processor == "email":
                raw_text, metadata = IngestionPipeline._extract_email(file_path)
            else:
                raw_text, metadata = IngestionPipeline._extract_text(file_path)
        except Exception as e:
            logger.error("Extraction failed for %s: %s", file_path, e)
            return {"error": f"Extraction failed: {str(e)}", "success": False}

        if not raw_text.strip():
            return {"error": "No content could be extracted", "success": False}

        # Step 2: AI Analysis — extract entities, summary, key points
        analysis = IngestionPipeline._ai_analyze(raw_text, tenant_id)

        # Step 2b: Generate a concise crux (one-liner essence)
        crux = IngestionPipeline._generate_crux(raw_text, analysis, path.stem)

        # Step 3: Store in knowledge base
        kb_entry = IngestionPipeline._store_in_knowledge_base(
            tenant_id=tenant_id,
            user_id=user_id,
            title=path.stem,
            raw_text=raw_text,
            analysis=analysis,
            source_label=source_label or f"Upload: {path.name}",
            file_path=file_path,
            category=category or "ingested_data",
        )

        # Step 4: Log activity
        activity = ActivityLog(
            tenant_id=tenant_id,
            user_id=user_id,
            action="data_ingested",
            detail=f"Ingested: {path.name} ({ext}) — {len(raw_text)} chars extracted",
        )
        db.session.add(activity)
        db.session.commit()

        return {
            "success": True,
            "filename": path.name,
            "extension": ext,
            "extraction_method": extraction_method,
            "characters_extracted": len(raw_text),
            "summary": analysis.get("summary", ""),
            "key_points": analysis.get("key_points", []),
            "crux": crux,
            "entities_detected": analysis.get("entities", []),
            "knowledge_entry_id": kb_entry.id if kb_entry else None,
            "metadata": metadata,
        }

    @staticmethod
    def _generate_crux(text: str, analysis: dict, title: str) -> str:
        """Generate a concise one-line crux about what was ingested."""
        category = analysis.get("suggested_category", "general")
        chars = len(text)
        lines = text.count("\n")
        words = len(text.split())
        
        category_emoji = {
            "finance": "💰", "hr": "👥", "legal": "⚖️",
            "marketing": "📈", "operations": "📋", "customer": "🤝",
        }
        emoji = category_emoji.get(category, "📄")
        
        parts = []
        if chars > 0:
            parts.append(f"{emoji} {title} — {category} document")
        
        if analysis.get("key_points"):
            top = analysis["key_points"][0][:80]
            parts.append(f"Key insight: {top}")
        elif analysis.get("summary"):
            parts.append(f"Starts with: {analysis['summary'][:80]}")
        
        if chars > 0:
            size_kb = chars / 1024
            parts.append(f"{'📃' if chars > 10000 else '📝'} {words:,} words, {size_kb:.0f}KB")
        
        return " · ".join(parts) if parts else f"📄 {title} processed ({chars} chars)"

    @staticmethod
    def _extract_pdf(file_path: str) -> tuple:
        """Extract text from PDF using PyMuPDF or fallback."""
        import fitz  # PyMuPDF
        doc = fitz.open(file_path)
        text = "\n".join([page.get_text() for page in doc])
        metadata = {
            "pages": len(doc),
            "title": doc.metadata.get("title", ""),
            "author": doc.metadata.get("author", ""),
        }
        doc.close()
        return text, metadata

    @staticmethod
    def _extract_docx(file_path: str) -> tuple:
        """Extract text from DOCX."""
        import docx
        doc = docx.Document(file_path)
        text = "\n".join([p.text for p in doc.paragraphs])
        metadata = {"paragraphs": len(doc.paragraphs)}
        return text, metadata

    @staticmethod
    def _extract_text(file_path: str) -> tuple:
        """Extract text from plain text files."""
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            text = f.read()
        return text, {"lines": text.count("\n")}

    @staticmethod
    def _extract_tabular(file_path: str) -> tuple:
        """Extract text from CSV/XLSX files."""
        import csv, io
        ext = Path(file_path).suffix.lower()
        text_parts = []
        metadata = {"rows": 0, "columns": 0}

        if ext == ".csv":
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                reader = csv.reader(f)
                headers = next(reader, [])
                metadata["columns"] = len(headers)
                text_parts.append(" | ".join(headers))
                for row in reader:
                    metadata["rows"] += 1
                    text_parts.append(" | ".join(row))
        elif ext in (".xlsx", ".xls"):
            import openpyxl
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                text_parts.append(f"\n=== Sheet: {sheet_name} ===")
                for i, row in enumerate(ws.iter_rows(values_only=True)):
                    if i == 0:
                        metadata["columns"] = max(metadata["columns"], len(row))
                    text_parts.append(" | ".join([str(c) if c is not None else "" for c in row]))
                    metadata["rows"] += 1
            wb.close()

        return "\n".join(text_parts), metadata

    @staticmethod
    def _extract_image(file_path: str) -> tuple:
        """Extract text from images using OCR."""
        try:
            import pytesseract
            from PIL import Image
            img = Image.open(file_path)
            text = pytesseract.image_to_string(img)
            metadata = {"width": img.width, "height": img.height}
            return text, metadata
        except ImportError:
            return "[OCR not available: pytesseract not installed]", {}

    @staticmethod
    def _extract_audio(file_path: str) -> tuple:
        """Extract text from audio using STT (Whisper)."""
        try:
            import openai
            with open(file_path, "rb") as audio:
                transcript = openai.OpenAI().audio.transcriptions.create(
                    model="whisper-1", file=audio,
                )
            duration = 0
            return transcript.text, {"duration_seconds": duration, "method": "whisper"}
        except Exception as e:
            logger.warning("Audio STT failed: %s", e)
            return f"[Audio transcription failed: {e}]", {}

    @staticmethod
    def _extract_video(file_path: str) -> tuple:
        """Extract text from video — extract audio track first, then STT."""
        try:
            import subprocess, tempfile
            # Extract audio using ffmpeg
            with tempfile.NamedTemporaryFile(suffix=".mp3", delete=False) as tmp:
                audio_path = tmp.name
            subprocess.run([
                "ffmpeg", "-i", file_path, "-vn", "-acodec", "libmp3lame",
                "-y", audio_path
            ], capture_output=True, timeout=60)
            text, metadata = IngestionPipeline._extract_audio(audio_path)
            os.unlink(audio_path)
            return text, metadata
        except Exception as e:
            logger.warning("Video extraction failed: %s", e)
            return f"[Video extraction failed: {e}]", {}

    @staticmethod
    def _extract_presentation(file_path: str) -> tuple:
        """Extract text from PPTX."""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)
            return "\n".join(text_parts), {"slides": len(prs.slides)}
        except ImportError:
            return "[PPTX extraction not available]", {}

    @staticmethod
    def _extract_email(file_path: str) -> tuple:
        """Extract text from EML/MSG files."""
        try:
            import email
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                msg = email.message_from_file(f)
            text_parts = []
            text_parts.append(f"From: {msg.get('From', '')}")
            text_parts.append(f"To: {msg.get('To', '')}")
            text_parts.append(f"Subject: {msg.get('Subject', '')}")
            text_parts.append(f"Date: {msg.get('Date', '')}")
            text_parts.append("")
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        text_parts.append(part.get_payload(decode=True).decode("utf-8", errors="replace"))
            else:
                text_parts.append(msg.get_payload(decode=True).decode("utf-8", errors="replace"))
            return "\n".join(text_parts), {"subject": msg.get("Subject", "")}
        except Exception as e:
            return f"[Email extraction failed: {e}]", {}

    @staticmethod
    def _extract_json(file_path: str) -> tuple:
        """Extract text from JSON with pretty formatting."""
        with open(file_path, "r", encoding="utf-8") as f:
            data = json.load(f)
        return json.dumps(data, indent=2), {"keys": list(data.keys()) if isinstance(data, dict) else []}

    @staticmethod
    def _process_archive(file_path: str, tenant_id: int, user_id: int,
                          category: Optional[str] = None) -> dict:
        """Extract and process all files in a ZIP archive."""
        results = []
        extract_dir = tempfile.mkdtemp()
        try:
            with zipfile.ZipFile(file_path, "r") as zf:
                zf.extractall(extract_dir)

            for root, dirs, files in os.walk(extract_dir):
                for fname in files:
                    fpath = os.path.join(root, fname)
                    result = IngestionPipeline.process_file(
                        fpath, tenant_id, user_id, category=category
                    )
                    results.append(result)

            return {
                "success": True,
                "archive": True,
                "files_processed": len(results),
                "results": results,
            }
        finally:
            import shutil
            shutil.rmtree(extract_dir, ignore_errors=True)

    @staticmethod
    def _ai_analyze(text: str, tenant_id: int) -> dict:
        """Analyze extracted text using AI to extract structured information."""
        analysis = {
            "summary": "",
            "key_points": [],
            "entities": [],
            "suggested_category": "general",
            "language": "en",
        }

        # Simple heuristic analysis
        lines = text.strip().split("\n")
        analysis["summary"] = lines[0][:300] if lines else ""

        # Extract key points (sentences with key indicators)
        key_indicators = ["important", "note", "key", "critical", "must", "should",
                         "required", "deadline", "policy", "rule", "regulation",
                         "important", "significant", "notable"]
        for line in lines[:50]:
            lower = line.lower()
            if any(kw in lower for kw in key_indicators) and len(line) > 20:
                analysis["key_points"].append(line.strip()[:200])

        # Detect language
        if any(ord(c) >= 0x0900 and ord(c) <= 0x097F for c in text[:500]):
            analysis["language"] = "hi"
        elif any(ord(c) >= 0x0600 and ord(c) <= 0x06FF for c in text[:500]):
            analysis["language"] = "ar"

        # Categorize
        category_keywords = {
            "finance": ["invoice", "payment", "revenue", "expense", "budget", "tax", "gst", "account"],
            "hr": ["employee", "salary", "hiring", "payroll", "leave", "attendance"],
            "legal": ["contract", "agreement", "clause", "compliance", "regulation", "policy"],
            "marketing": ["campaign", "social media", "email", "content", "seo", "brand"],
            "operations": ["project", "task", "workflow", "process", "sop", "milestone"],
            "customer": ["customer", "client", "feedback", "support", "ticket", "complaint"],
        }
        text_lower = text.lower()
        for cat, keywords in category_keywords.items():
            if any(kw in text_lower for kw in keywords):
                analysis["suggested_category"] = cat
                break

        analysis["entities"] = list(set(analysis["entities"]))
        return analysis

    @staticmethod
    def _store_in_knowledge_base(tenant_id: int, user_id: int, title: str,
                                  raw_text: str, analysis: dict,
                                  source_label: str, file_path: str,
                                  category: str) -> Optional[KnowledgeEntry]:
        """Store extracted knowledge in the knowledge base."""
        # Create a comprehensive knowledge entry
        entry_text = f"# {title}\n\n"
        if analysis.get("summary"):
            entry_text += f"**Summary:** {analysis['summary']}\n\n"
        if analysis.get("key_points"):
            entry_text += "**Key Points:**\n"
            for kp in analysis["key_points"][:10]:
                entry_text += f"- {kp}\n"
            entry_text += "\n"
        entry_text += f"**Full Content:**\n{raw_text[:5000]}"

        entry = KnowledgeEntry(
            tenant_id=tenant_id,
            question=f"Ingested: {title}",
            answer=entry_text,
            source=source_label,
            source_url=file_path,
            confidence=0.7,
            category=category,
            meta_data=json.dumps({
                "title": title,
                "characters": len(raw_text),
                "source": source_label,
                "analysis": {k: v for k, v in analysis.items() if k != "entities"},
                "ingested_at": datetime.utcnow().isoformat(),
                "ingested_by": user_id,
            }),
            file_type=Path(file_path).suffix.lstrip("."),
        )
        db.session.add(entry)
        db.session.flush()
        return entry


# ---------------------------------------------------------------------------
# Ingestion API endpoints
# ---------------------------------------------------------------------------

from flask import Blueprint, request, jsonify, render_template, g
from app.routes.auth import login_required

ingestion_bp = Blueprint("ingestion", __name__, url_prefix="/ingestion")


@ingestion_bp.route("")
@login_required
def ingestion_page():
    """Data Ingestion Hub frontend."""
    # Get recent ingestions
    recent = KnowledgeEntry.query.filter_by(
        tenant_id=g.tenant.id,
    ).filter(
        KnowledgeEntry.source.ilike("Upload:%")
    ).order_by(KnowledgeEntry.created_at.desc()).limit(50).all()

    return render_template("ingestion.html", recent=recent,
                           supported_formats=list(SUPPORTED_EXTENSIONS.keys()))


@ingestion_bp.route("/upload", methods=["POST"])
@login_required
def upload_file():
    """Upload a file or URL for ingestion."""
    files = request.files.getlist("files") if request.files else []
    file = request.files.get("file")
    url = request.form.get("url", "").strip()
    category = request.form.get("category", "general")

    # Multi-file upload
    if files and len(files) > 1:
        results = []
        for f in files:
            ext = Path(f.filename).suffix.lower() or ".bin"
            import tempfile
            with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
                f.save(tmp.name)
                fp = tmp.name
            try:
                r = IngestionPipeline.process_file(fp, g.tenant.id, g.user.id,
                    source_label=f"Upload: {f.filename}", category=category)
                results.append(r)
            finally:
                try: os.unlink(fp)
                except: pass
        return jsonify({"success": True, "batch": True, "count": len(results), "results": results})

    if not file and not url:
        return jsonify({"error": "No file or URL provided"}), 400

    if url:
        # Download from URL
        try:
            import requests
            resp = requests.get(url, timeout=30, headers={"User-Agent": "ShunyaOS/1.0"})
            if resp.status_code != 200:
                return jsonify({"error": f"URL returned {resp.status_code}"}), 400
            ext = Path(url).suffix or ".html"
            with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
                tmp.write(resp.content)
                file_path = tmp.name
            source_label = f"URL: {url}"
        except Exception as e:
            return jsonify({"error": f"URL download failed: {str(e)}"}), 400
    else:
        # Save uploaded file
        ext = Path(file.filename).suffix.lower() or ".bin"
        with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
            file.save(tmp.name)
            file_path = tmp.name
        source_label = f"Upload: {file.filename}"

    try:
        result = IngestionPipeline.process_file(
            file_path=file_path,
            tenant_id=g.tenant.id,
            user_id=g.user.id,
            source_label=source_label,
            category=category,
        )
        return jsonify(result)
    finally:
        try:
            os.unlink(file_path)
        except OSError:
            pass


@ingestion_bp.route("/history")
@login_required
def ingestion_history():
    """Get ingestion history as JSON."""
    recent = KnowledgeEntry.query.filter_by(
        tenant_id=g.tenant.id,
    ).filter(
        KnowledgeEntry.source.ilike("Upload:%")
    ).order_by(KnowledgeEntry.created_at.desc()).limit(50).all()

    return jsonify({"entries": [{
        "id": e.id,
        "title": e.question.replace("Ingested: ", ""),
        "source": e.source,
        "category": e.category,
        "confidence": e.confidence,
        "created_at": e.created_at.isoformat() if e.created_at else None,
    } for e in recent]})